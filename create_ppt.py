from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_ppt():
    prs = Presentation()

    def set_slide_background(slide, r, g, b):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)

    def add_title_slide():
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, 255, 255, 255)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "LocalFix"
        title.text_frame.paragraphs[0].font.size = Pt(60)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(79, 70, 229) # Indigo/Primary
        
        subtitle.text = "Home Service Booking Platform\nProfessional & Reliable Solutions"
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)

    def add_overview_slide():
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "Project Overview"
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.text = "Connecting users with elite home service providers."
        
        p = tf.add_paragraph()
        p.text = "Key Objectives:"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = "• Simplify home maintenance & repairs"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Provide certified & background-checked professionals"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "• Seamless role-based user experience (User, Provider, Admin)"
        p.level = 1

    def add_screenshot_slide(title_text, image_path):
        slide_layout = prs.slide_layouts[5] # Blank with Title
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title_text
        
        if os.path.exists(image_path):
            # Centering the image
            # SLIDE_WIDTH = 10 inches, SLIDE_HEIGHT = 7.5 inches
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            slide.shapes.add_picture(image_path, left, top, width=width)
        else:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
            tf = txBox.text_frame
            tf.text = f"[Image not found: {os.path.basename(image_path)}]"
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    def add_features_slide():
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Key Features"
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        features = [
            "Premium Modern UI with Glassmorphism",
            "Role-Based Dashboards",
            "Real-time Service Tracking",
            "Secure Managed Bookings",
            "Review & Rating System",
            "Interactive Admin Control Panel"
        ]
        
        for feature in features:
            p = tf.add_paragraph()
            p.text = feature
            p.level = 0

    def add_tech_stack_slide():
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Technology Stack"
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        stack = {
            "Frontend": "React, TypeScript, Vite",
            "Styling": "Tailwind CSS, Framer Motion",
            "State Management": "Zustand",
            "Backend": "Node.js, Express",
            "Database": "MongoDB"
        }
        
        for key, value in stack.items():
            p = tf.add_paragraph()
            p.text = f"{key}: {value}"
            p.level = 0

    # Main Execution
    add_title_slide()
    add_overview_slide()
    add_features_slide()
    add_screenshot_slide("Landing Page interface", "c:\\Users\\Home1\\Downloads\\projectlocalfix\\landing_page.png")
    add_screenshot_slide("User Dashboard", "c:\\Users\\Home1\\Downloads\\projectlocalfix\\user_dashboard.png")
    add_tech_stack_slide()
    
    output_path = "c:\\Users\\Home1\\Downloads\\projectlocalfix\\LocalFix_Project_Presentation.pptx"
    prs.save(output_path)
    print(f"PPT created successfully at {output_path}")

if __name__ == "__main__":
    create_ppt()
